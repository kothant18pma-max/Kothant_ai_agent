import os
import threading
import tempfile
import google.generativeai as genai
from huggingface_hub import InferenceClient
from openai import OpenAI
from flask import Flask
from pptx import Presentation # PowerPoint အတွက် (pip install python-pptx)
from openpyxl import Workbook # Excel အတွက် (pip install openpyxl)
from docx import Document # Word အတွက် (pip install python-docx)
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import ApplicationBuilder, ContextTypes, MessageHandler, filters, CallbackQueryHandler

# --- API Clients Setup ---
grok_client = OpenAI(api_key=os.environ.get("XAI_API_KEY"), base_url="https://api.x.ai/v1") if os.environ.get("XAI_API_KEY") else None
ds_client = OpenAI(api_key=os.environ.get("DEEPSEEK_API_KEY"), base_url="https://api.deepseek.com") if os.environ.get("DEEPSEEK_API_KEY") else None
hf_client = InferenceClient(model="mistralai/Mistral-7B-Instruct-v0.3", token=os.environ.get("HF_TOKEN")) if os.environ.get("HF_TOKEN") else None

# --- Flask Server ---
flask_app = Flask(__name__)
@flask_app.route('/')
def home(): return "AI Agent is Active!"

def run_flask():
    port = int(os.environ.get("PORT", 10000))
    flask_app.run(host='0.0.0.0', port=port)

# --- Utilities ---

async def send_split_message(update_or_query, text):
    MAX_LENGTH = 4000
    for i in range(0, len(text), MAX_LENGTH):
        part = text[i:i+MAX_LENGTH]
        if isinstance(update_or_query, Update):
            await update_or_query.message.reply_text(part)
        else:
            await update_or_query.message.chat.send_message(part)

def create_pptx(content_text, output_path):
    prs = Presentation()
    lines = content_text.split('\n')
    
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AI Generated Presentation"
    slide.placeholders[1].text = "Summarized by AI Bot"

    for i in range(0, len(lines), 5):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Key Points (Part {i//5 + 1})"
        slide.placeholders[1].text = "\n".join(lines[i:i+5])
            
    prs.save(output_path)

def create_excel(content_text, output_path):
    """AI ပေးသော CSV data ကို Excel အဖြစ် ပြောင်းပေးရန်"""
    wb = Workbook()
    ws = wb.active
    ws.title = "Extracted Data"
    
    lines = content_text.strip().split('\n')
    for line in lines:
        if not line.strip(): continue # တိုက်ရိုက် line များကို ချန်လှပ်
        # Comma (,) ဖြင့် ခွဲထားသော data များကို Cell အလိုက် ထည့်သည်
        cells = [cell.strip() for cell in line.split(',')]
        ws.append(cells)
        
    wb.save(output_path)

def create_word(content_text, output_path):
    """AI ပေးသော စာသားကို Word အဖြစ် ပြောင်းပေးရန်"""
    doc = Document()
    doc.add_heading('Extracted Document', 0)
    
    lines = content_text.split('\n')
    for line in lines:
        if line.strip():
            doc.add_paragraph(line)
            
    doc.save(output_path)

# --- AI Core Functions ---

def get_gemini_res(prompt, file_path=None):
    key = os.environ.get("GOOGLE_API_KEY")
    if not key: return "Error: Gemini Key missing."
    genai.configure(api_key=key)
    model = genai.GenerativeModel('gemini-2.5-flash') 
    if file_path:
        # Gemini သည် PDF နှင့် Image နှစ်မျိုးလုံးကို ဖတ်နိုင်သည်
        up_file = genai.upload_file(path=file_path)
        return model.generate_content([prompt, up_file]).text
    return model.generate_content(prompt).text

def get_ai_chat(text):
    for func in [lambda: ds_client.chat.completions.create(model="deepseek-chat", messages=[{"role": "user", "content": text}]).choices[0].message.content if ds_client else None,
                 lambda: grok_client.chat.completions.create(model="grok-beta", messages=[{"role": "user", "content": text}]).choices[0].message.content if grok_client else None]:
        try:
            res = func()
            if res: return res
        except: continue
    return get_gemini_res(f"Reply in Myanmar: {text}")

# --- Handlers ---

async def handle_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    res = get_ai_chat(update.message.text)
    await send_split_message(update, res)

async def handle_media(update: Update, context: ContextTypes.DEFAULT_TYPE):
    # ဖိုင်အမျိုးအစား စစ်ဆေးခြင်း (JPEG နှင့် PDF အတွက် အထူးပြုလုပ်ထားသည်)
    if update.message.document:
        file_obj = await update.message.document.get_file()
        # File extension ကို အလိုအလျောက် ယူသည် (ဥပမာ .pdf, .jpg)
        original_name = update.message.document.file_name
        _, ext = os.path.splitext(original_name)
        suffix = ext if ext else ".pdf" # extension မရှိရင် PDF အဖြစ် သတ်မှတ်
    elif update.message.photo:
        file_obj = await update.message.photo[-1].get_file()
        suffix = ".jpg" # Telegram photos are jpg
    else:
        return

    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as t:
        await file_obj.download_to_drive(t.name)
        context.user_data['current_file'] = t.name
    
    # Word နှင့် Excel ခလုတ်များ ပါဝင်သော Menu
    kb = [
        [InlineKeyboardButton("🔍 OCR", callback_data='ocr'), InlineKeyboardButton("📝 Summary", callback_data='sum')],
        [InlineKeyboardButton("📊 Create PPTX", callback_data='pptx'), InlineKeyboardButton("📈 Create Excel", callback_data='excel')],
        [InlineKeyboardButton("📄 Create Word", callback_data='docx')]
    ]
    await update.message.reply_text("📁 ဖိုင်ရပါပြီ။ ဘာလုပ်ရမလဲ?", reply_markup=InlineKeyboardMarkup(kb))

async def button_click(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    file_path = context.user_data.get('current_file')
    if not file_path: return

    cmd = query.data
    
    if cmd == 'pptx':
        await query.edit_message_text("⚙️ PowerPoint ဖန်တီးနေပါသည်...")
        prompt = "Extract key facts from this file and list them line by line."
        res = get_gemini_res(prompt, file_path)
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as t:
            create_pptx(res, t.name)
            await context.bot.send_document(chat_id=query.message.chat_id, document=open(t.name, 'rb'), filename="AI_Presentation.pptx")
            
    elif cmd == 'excel':
        await query.edit_message_text("⚙️ Excel ဖိုင်ဖန်တီးနေပါသည်... (Data များ ပြန်လည်စီစစ်နေပါသည်)")
        # JPEG/PDF ထဲမှ Data များကို CSV format နှင့် တိုက်ရိုက်ထုတ်ပေးရန် မေးမြန်းခြင်း
        prompt = (
            "Analyze the provided document or image. "
            "Extract all data structures or tables. "
            "Output ONLY valid CSV format (comma-separated values). "
            "Do not add any explanation or markdown code blocks. "
            "If there are no tables, list the key information in rows."
        )
        res = get_gemini_res(prompt, file_path)
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as t:
            create_excel(res, t.name)
            await context.bot.send_document(chat_id=query.message.chat_id, document=open(t.name, 'rb'), filename="AI_Data.xlsx")
            
    elif cmd == 'docx':
        await query.edit_message_text("⚙️ Word ဖိုင်ဖန်တီးနေပါသည်... (OCR ဆောင်ရွက်နေပါသည်)")
        # JPEG/PDF ထဲမှ စာသားများကို အပြည့်အစုံ ထုတ်ပေးရန် မေးမြန်းခြင်း
        prompt = "Perform OCR on this file. Extract all text content clearly. Maintain original structure if possible."
        res = get_gemini_res(prompt, file_path)
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as t:
            create_word(res, t.name)
            await context.bot.send_document(chat_id=query.message.chat_id, document=open(t.name, 'rb'), filename="AI_Document.docx")
            
    else:
        prompt = "Extract all text clearly." if cmd == 'ocr' else "Summarize this clearly in Myanmar."
        await query.edit_message_text(f"⚙️ {cmd.upper()} လုပ်ဆောင်နေပါသည်...")
        res = get_gemini_res(prompt, file_path)
        await send_split_message(query, res)

if __name__ == '__main__':
    threading.Thread(target=run_flask, daemon=True).start()
    token = os.environ.get("TELEGRAM_TOKEN")
    if token:
        app = ApplicationBuilder().token(token).build()
        app.add_handler(MessageHandler(filters.TEXT & (~filters.COMMAND), handle_text))
        app.add_handler(MessageHandler(filters.Document.ALL | filters.PHOTO, handle_media))
        app.add_handler(CallbackQueryHandler(button_click))
        app.run_polling(drop_pending_updates=True)
